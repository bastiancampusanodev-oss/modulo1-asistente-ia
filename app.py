# ============================================================
# DEMO CLASE 3 / MÓDULO 1
# Asistente IA con FastAPI + LangChain + OpenAI + Chroma
# ============================================================
#
#
# 1. Primero importamos las librerías.
# 2. Después cargamos variables de entorno y configuraciones.
# 3. Luego definimos un contexto fijo del curso.
# 4. Creamos una memoria simple de conversación en RAM.
# 5. Levantamos la app FastAPI y definimos los modelos de entrada/salida.
# 6. Creamos funciones auxiliares para memoria, documentos, vector store y RAG.
# 7. Creamos un router/clasificador que decide cómo responder.
# 8. Creamos prompts para responder con RAG o recomendar arquitectura.
# 9. Finalmente definimos los endpoints de la API.
#
# ============================================================

# 1) IMPORTS BÁSICOS DE PYTHON
# os permite leer variables de entorno, por ejemplo OPENAI_MODEL.
import os
import shutil
from pathlib import Path
from typing import Dict, List


# 2) LIBRERÍAS DE BACKEND/API Y VALIDACIÓN
# dotenv carga variables desde el archivo .env.
from dotenv import load_dotenv
from fastapi import FastAPI
from fastapi.responses import HTMLResponse
from pydantic import BaseModel
from pptx import Presentation
from openai import OpenAI


# 3) LIBRERÍAS DE LANGCHAIN / RAG
# ChatOpenAI conecta con el modelo conversacional.
# OpenAIEmbeddings convierte texto en vectores para búsqueda semántica.
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_chroma import Chroma
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate

# 4) CARGA DE CONFIGURACIÓN
# Busca un archivo .env y carga claves/modelos configurados ahí.
load_dotenv()

# Carpeta donde dejamos los materiales del curso: PPTX, PDF, TXT o MD.
DATA_DIR = Path("data")

# Carpeta local donde Chroma guardará el índice vectorial.
# Si ya existe, se reutiliza para no reconstruir embeddings cada vez.
DB_DIR = Path("/tmp/chroma_modulo1")
# Modelo principal para clasificar y responder.
MODEL_NAME = os.getenv("OPENAI_MODEL", "gpt-4.1-mini")
# Modelo usado cuando activamos búsqueda web.
WEB_MODEL_NAME = os.getenv("OPENAI_WEB_MODEL", MODEL_NAME)
# Modelo de embeddings: transforma texto en vectores.
EMBEDDING_MODEL = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")

# 5) CONTEXTO FIJO DEL CURSO
# Este texto no viene del usuario ni del RAG: es conocimiento estable
# que la app siempre tendrá disponible para responder datos administrativos.
COURSE_CONTEXT = """
Datos fijos del curso:
- Profesor / relator: Bastián Campusano.
- Módulo: Módulo 1 · APIs de IA, LangChain y RAG.
- Programa: Diplomado Ejecutivo en Confianza Digital, IA y Automatización Empresarial.
- Institución: USACH.
- Modalidad: clase virtual.
- Objetivo del módulo: diseñar soluciones IA simples, seguras, medibles y defendibles.
"""

# Memoria simple en RAM para demo.
# Se borra si reinicias uvicorn. Es perfecta para explicar el concepto en clase.
CONVERSATION_MEMORY: Dict[str, List[Dict[str, str]]] = {}
MAX_MEMORY_MESSAGES = 8

# 6) CREACIÓN DE LA APLICACIÓN FASTAPI
# Aquí nace el backend. FastAPI expone endpoints como /api/ask o /api/health.
app = FastAPI(title="Asistente Módulo 1")

# Cliente directo de OpenAI. Lo usamos para la herramienta de web_search.
client = OpenAI()


# 7) MODELOS DE DATOS DE LA API
# Pydantic define y valida la estructura de los datos que entran y salen.
# AskRequest representa lo que manda el frontend cuando el usuario pregunta.
class AskRequest(BaseModel):
    question: str
    session_id: str = "default"


# AskResponse representa la respuesta estándar del backend al frontend.
class AskResponse(BaseModel):
    answer: str
    mode: str
    sources: List[str]


# ClearMemoryRequest permite limpiar la memoria de una sesión específica.
class ClearMemoryRequest(BaseModel):
    session_id: str = "default"


# 8) MEMORIA CONVERSACIONAL SIMPLE
# Recupera el historial asociado a una sesión.
# En producción esto podría guardarse en Redis, Postgres u otra base de datos.
def get_memory(session_id: str) -> List[Dict[str, str]]:
    return CONVERSATION_MEMORY.get(session_id, [])


# Guarda un turno completo: mensaje del usuario + respuesta del asistente.
# Se limita a MAX_MEMORY_MESSAGES para que el prompt no crezca indefinidamente.
def save_turn(session_id: str, user_message: str, assistant_message: str) -> None:
    CONVERSATION_MEMORY.setdefault(session_id, [])
    CONVERSATION_MEMORY[session_id].append({"role": "user", "content": user_message})
    CONVERSATION_MEMORY[session_id].append({"role": "assistant", "content": assistant_message})
    CONVERSATION_MEMORY[session_id] = CONVERSATION_MEMORY[session_id][-MAX_MEMORY_MESSAGES:]


# Convierte la memoria en texto plano para inyectarla en los prompts.
# Esto permite preguntas como "¿y en ese caso?" usando contexto previo.
def format_memory(session_id: str) -> str:
    memory = get_memory(session_id)
    if not memory:
        return "No hay historial previo en esta sesión."

    lines = []
    for item in memory[-MAX_MEMORY_MESSAGES:]:
        role = "Usuario" if item["role"] == "user" else "Asistente"
        content = item["content"].replace("\n", " ").strip()
        lines.append(f"{role}: {content}")
    return "\n".join(lines)


# Limpia la memoria de una sesión.
# Útil para mostrar en clase qué pasa con y sin historial.
def clear_session_memory(session_id: str) -> None:
    CONVERSATION_MEMORY.pop(session_id, None)


# 9) CARGA DE DOCUMENTOS DEL CURSO
# Esta función lee una presentación PPTX y extrae texto slide por slide.
# Cada slide se transforma en un Document de LangChain con metadata.
def load_pptx_documents(file_path: Path) -> List[Document]:
    presentation = Presentation(str(file_path))
    docs = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        slide_text = "\n".join(texts).strip()
        if slide_text:
            docs.append(
                Document(
                    page_content=slide_text,
                    metadata={"source": file_path.name, "slide": slide_index, "type": "pptx"},
                )
            )
    return docs


# Carga todos los documentos disponibles en la carpeta data/.
# Soporta PPTX, PDF, TXT y MD.
# Resultado: una lista de Documents lista para chunking y embeddings.
def load_documents() -> List[Document]:
    docs = []

    for file_path in DATA_DIR.glob("*"):
        suffix = file_path.suffix.lower()

        if suffix == ".pptx":
            docs.extend(load_pptx_documents(file_path))

        elif suffix == ".pdf":
            loader = PyPDFLoader(str(file_path))
            loaded = loader.load()
            for doc in loaded:
                doc.metadata["source"] = file_path.name
                doc.metadata["type"] = "pdf"
            docs.extend(loaded)

        elif suffix in [".txt", ".md"]:
            loader = TextLoader(str(file_path), encoding="utf-8")
            loaded = loader.load()
            for doc in loaded:
                doc.metadata["source"] = file_path.name
                doc.metadata["type"] = suffix.replace(".", "")
            docs.extend(loaded)

    return docs


# 10) VECTOR STORE / ÍNDICE RAG
# Esta es una de las partes centrales de la demo.
# - Crea embeddings.
# - Si existe un índice Chroma previo, lo reutiliza.
# - Si no existe, carga documentos, los divide en chunks y crea el índice.
def get_vectorstore(force_rebuild: bool = False):
    embeddings = OpenAIEmbeddings(model=EMBEDDING_MODEL)

    if force_rebuild and DB_DIR.exists():
        shutil.rmtree(DB_DIR)

    if DB_DIR.exists():
        return Chroma(persist_directory=str(DB_DIR), embedding_function=embeddings)

    documents = load_documents()
    if not documents:
        raise RuntimeError("No encontré documentos en /data. Agrega PPTX, PDF, TXT o MD del Módulo 1.")

    splitter = RecursiveCharacterTextSplitter(chunk_size=900, chunk_overlap=120)
    chunks = splitter.split_documents(documents)

    return Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=str(DB_DIR),
    )


# Convierte los documentos recuperados por el retriever en texto para el prompt.
# También agrega fuente y ubicación para que la respuesta sea más trazable.
def format_docs(docs: List[Document]) -> str:
    if not docs:
        return "No se recuperaron fragmentos desde los materiales del módulo."

    parts = []
    for i, doc in enumerate(docs, start=1):
        source = doc.metadata.get("source", "fuente_desconocida")
        slide = doc.metadata.get("slide")
        page = doc.metadata.get("page")

        location = ""
        if slide:
            location = f"slide {slide}"
        elif page is not None:
            location = f"página {page}"

        parts.append(f"[Material {i}: {source} {location}]\n{doc.page_content}")

    return "\n\n".join(parts)


# Extrae nombres de fuentes recuperadas para mostrarlas en la interfaz.
# Esto ayuda a explicar que RAG no solo responde: también puede mostrar evidencia.
def extract_sources(docs: List[Document]) -> List[str]:
    sources = []

    for doc in docs:
        source = doc.metadata.get("source", "fuente_desconocida")
        slide = doc.metadata.get("slide")
        page = doc.metadata.get("page")

        if slide:
            label = f"{source} · slide {slide}"
        elif page is not None:
            label = f"{source} · página {page}"
        else:
            label = source

        if label not in sources:
            sources.append(label)

    return sources[:6]


# 11) ROUTER / CLASIFICADOR DE PREGUNTAS
# Esta función decide qué camino seguirá la pregunta:
# - MODULO: responder con contexto del curso + RAG.
# - WEB: usar búsqueda web.
# - DECISION: recomendar arquitectura.
# - MIXTA: combinar web + módulo.
#
# Este router es una forma simple de mostrar "orquestación".
def classify_question(question: str, memory_context: str) -> str:
    model = ChatOpenAI(model=MODEL_NAME, temperature=0)

    prompt = ChatPromptTemplate.from_messages([
        (
            "system",
            """
Clasifica la pregunta en una sola categoría.

Categorías posibles:
MODULO: pregunta sobre clases, conceptos del módulo, APIs, LangChain, RAG, prompts, evaluación o contenido cargado.
WEB: pregunta técnica general o actual que puede requerir internet, ejemplos externos, herramientas reales, documentación vigente o nombres de servicios.
DECISION: pregunta donde se debe elegir arquitectura: API simple, chain, RAG, tool, agente, dashboard sin IA o no usar IA.
MIXTA: combina contenido del módulo con información externa.

Usa el historial solo para entender referencias como "eso", "lo anterior" o "ese caso".
Responde solo con una palabra: MODULO, WEB, DECISION o MIXTA.
""",
        ),
        (
            "user",
            """
Pregunta actual:
{question}

Historial reciente de la sesión:
{memory_context}
""",
        ),
    ])

    result = (prompt | model).invoke({"question": question, "memory_context": memory_context})
    category = result.content.strip().upper()

    if category not in ["MODULO", "WEB", "DECISION", "MIXTA"]:
        return "MODULO"
    return category


# 12) BÚSQUEDA WEB
# Se usa cuando el router decide que la pregunta necesita información externa
# o actualizada, por ejemplo nombres de herramientas o ejemplos vigentes.
def run_web_search(question: str, memory_context: str) -> str:
    response = client.responses.create(
        model=WEB_MODEL_NAME,
        tools=[{"type": "web_search"}],
        input=f"""
Busca información actualizada si es necesario y responde en español claro.

Pregunta:
{question}

Historial reciente:
{memory_context}

Reglas de respuesta:
- Máximo 140 palabras.
- Máximo 5 bullets.
- No incluyas código salvo que el usuario lo pida explícitamente.
- Si piden ejemplos, entrega solo 3 a 5 ejemplos.
- Cierra con una frase práctica.
""",
    )
    return response.output_text


# 13) RESPUESTA CON CONTEXTO DEL MÓDULO + RAG
# Aquí generamos el prompt principal para responder preguntas del curso.
# El modelo recibe:
# - pregunta del usuario,
# - contexto fijo,
# - memoria de sesión,
# - fragmentos recuperados desde los materiales.
def answer_with_module(question: str, module_context: str, memory_context: str) -> str:
    model = ChatOpenAI(model=MODEL_NAME, temperature=0.2)

    prompt = ChatPromptTemplate.from_messages([
        (
            "system",
            """
Eres un asistente académico del Módulo 1: APIs de IA, LangChain y RAG.

Tienes tres fuentes:
1. Contexto fijo del sistema: datos administrativos estables del curso.
2. Memoria de sesión: historial reciente de esta conversación.
3. RAG: fragmentos recuperados desde los materiales del módulo.

Prioridad:
- Usa el contexto fijo para datos administrativos del curso.
- Usa la memoria para entender continuidad conversacional.
- Usa RAG para responder contenidos de clase.
- Si el contexto no alcanza, dilo explícitamente.

Reglas:
1. Responde en español claro, profesional y didáctico.
2. No inventes contenidos del curso.
3. Responde en máximo 140 palabras.
4. Usa máximo 5 bullets.
5. No incluyas código salvo que el usuario lo pida explícitamente.
6. Si la pregunta es conceptual, responde con: definición breve, para qué sirve y ejemplo corto.
7. Cierra con una frase práctica para que el estudiante sepa qué decisión tomar.
""",
        ),
        (
            "user",
            """
Pregunta:
{question}

Contexto fijo del sistema:
{course_context}

Memoria de sesión:
{memory_context}

Contexto recuperado desde materiales del módulo:
{module_context}
""",
        ),
    ])

    return (prompt | model).invoke({
        "question": question,
        "course_context": COURSE_CONTEXT,
        "memory_context": memory_context,
        "module_context": module_context,
    }).content


# 14) RESPUESTA PARA DECISIONES DE ARQUITECTURA
# Esta función no solo "responde": recomienda un patrón técnico.
# Sirve para mostrar cuándo usar API simple, chain, RAG, tool, agente o no IA.
def answer_architecture_decision(question: str, module_context: str, memory_context: str) -> str:
    model = ChatOpenAI(model=MODEL_NAME, temperature=0.2)

    prompt = ChatPromptTemplate.from_messages([
        (
            "system",
            """
Eres un arquitecto de soluciones IA para procesos de negocio.

Tienes tres fuentes:
1. Contexto fijo del sistema: datos del curso.
2. Memoria de sesión: historial reciente para entender continuidad.
3. RAG: materiales del módulo.

Tu tarea es recomendar el patrón más simple y correcto:
- No usar IA
- Dashboard / BI tradicional
- API simple a LLM
- Chain con structured output
- RAG
- Tool
- Agente controlado
- Workflow con revisión humana

Criterio:
1. Si es cálculo, validación simple, dashboard o alerta por umbral, recomienda no usar IA o usar reglas/BI.
2. Si es transformar texto sin datos externos, recomienda API simple o chain.
3. Si requiere salida estructurada, recomienda chain + parser/schema.
4. Si depende de documentos, políticas o evidencia, recomienda RAG.
5. Si necesita consultar datos vivos en sistemas, recomienda tool.
6. Si necesita decidir entre varias rutas o herramientas, recomienda agente controlado.
7. Si hay riesgo legal, médico, financiero o reputacional, agrega revisión humana.

Entrega:
- Patrón recomendado.
- Por qué en 2 frases.
- Flujo propuesto en máximo 5 pasos.
- Qué NO usaría al inicio.
- Control mínimo.

Máximo 180 palabras.
Sin párrafos largos.
""",
        ),
        (
            "user",
            """
Caso o pregunta:
{question}

Contexto fijo del sistema:
{course_context}

Memoria de sesión:
{memory_context}

Contexto disponible desde el módulo:
{module_context}
""",
        ),
    ])

    return (prompt | model).invoke({
        "question": question,
        "course_context": COURSE_CONTEXT,
        "memory_context": memory_context,
        "module_context": module_context,
    }).content


# 15) ENDPOINTS / RUTAS DE LA APLICACIÓN
# Ruta principal: entrega el HTML de la interfaz visual.
@app.get("/", response_class=HTMLResponse)
def home():
    return HTMLResponse(INDEX_HTML)


# Health check: permite revisar si la app está viva, qué modelo usa,
# qué archivos hay en data/ y si existe el índice Chroma.
@app.get("/api/health")
def health():
    files = [p.name for p in DATA_DIR.glob("*")]
    return {
        "ok": True,
        "model": MODEL_NAME,
        "web_model": WEB_MODEL_NAME,
        "embedding_model": EMBEDDING_MODEL,
        "data_files": files,
        "db_exists": DB_DIR.exists(),
        "memory_sessions": list(CONVERSATION_MEMORY.keys()),
    }


# Reconstruye el índice vectorial.
# Útil cuando agregamos o cambiamos documentos en la carpeta data/.
@app.post("/api/rebuild")
def rebuild():
    get_vectorstore(force_rebuild=True)
    return {"ok": True, "message": "Índice reconstruido correctamente."}


# Limpia la memoria de conversación de la sesión actual.
@app.post("/api/memory/clear")
def clear_memory(payload: ClearMemoryRequest):
    clear_session_memory(payload.session_id)
    return {"ok": True, "message": "Memoria de sesión limpiada correctamente."}


# 16) ENDPOINT PRINCIPAL: /api/ask
# Este es el corazón de la demo.
# Recibe una pregunta y ejecuta el flujo completo:
# 1. Limpia/valida pregunta.
# 2. Recupera memoria.
# 3. Construye retriever.
# 4. Busca documentos relevantes.
# 5. Clasifica la pregunta con el router.
# 6. Decide si responder con RAG, web, decisión o modo mixto.
# 7. Guarda el turno en memoria.
# 8. Devuelve respuesta, modo y fuentes.
@app.post("/api/ask", response_model=AskResponse)
def ask(payload: AskRequest):
    question = payload.question.strip()
    session_id = payload.session_id.strip() or "default"

    if not question:
        return AskResponse(answer="Escribe una pregunta para poder ayudarte.", mode="VACIA", sources=[])

    try:
        # Recuperamos memoria reciente para entender continuidad conversacional.
        memory_context = format_memory(session_id)

        # Cargamos o creamos el vector store Chroma.
        vectorstore = get_vectorstore()
        # Creamos el retriever: buscará los 5 fragmentos más relevantes.
        retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

        retrieval_query = f"""
Pregunta actual:
{question}

Historial reciente:
{memory_context}
"""
        # Ejecutamos la búsqueda semántica sobre los materiales del módulo.
        docs = retriever.invoke(retrieval_query)
        # Formateamos los documentos recuperados para pasarlos al prompt.
        module_context = format_docs(docs)
        sources = extract_sources(docs)

        # El router decide qué tipo de respuesta corresponde.
        mode = classify_question(question, memory_context)

        # Camino 1: pregunta externa o actualizada -> web_search.
        if mode == "WEB":
            answer = run_web_search(question, memory_context)
            save_turn(session_id, question, answer)
            return AskResponse(answer=answer, mode=mode, sources=["OpenAI web_search"])

        # Camino 2: pregunta de arquitectura -> recomendación de patrón.
        if mode == "DECISION":
            answer = answer_architecture_decision(question, module_context, memory_context)
            save_turn(session_id, question, answer)
            return AskResponse(answer=answer, mode=mode, sources=sources + ["Memoria de sesión", "Contexto fijo"])

        # Camino 3: combina información web con contexto del módulo.
        if mode == "MIXTA":
            web_answer = run_web_search(question, memory_context)
            combined_question = f"""
Pregunta original:
{question}

Información web:
{web_answer}
"""
            answer = answer_with_module(combined_question, module_context, memory_context)
            save_turn(session_id, question, answer)
            return AskResponse(answer=answer, mode=mode, sources=sources + ["OpenAI web_search", "Memoria de sesión", "Contexto fijo"])

        # Camino 4 por defecto: respuesta académica usando contexto fijo + memoria + RAG.
        answer = answer_with_module(question, module_context, memory_context)
        save_turn(session_id, question, answer)
        return AskResponse(answer=answer, mode=mode, sources=sources + ["Memoria de sesión", "Contexto fijo"])

    except Exception as error:
        return AskResponse(
            answer=f"Error controlado en backend: {type(error).__name__}: {str(error)}",
            mode="ERROR",
            sources=[],
        )



# 17) FRONTEND HTML
# Desde aquí empieza la interfaz visual.
# Se deja prácticamente sin comentar para que en clase el foco esté en el backend.
INDEX_HTML = """
<!DOCTYPE html>
<html lang="es">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0"/>
  <title>Asistente Módulo 1</title>
  <style>
    :root {
      --bg: #080b12;
      --panel: rgba(15, 23, 42, 0.82);
      --panel2: rgba(30, 41, 59, 0.72);
      --border: rgba(148, 163, 184, 0.22);
      --text: #e5e7eb;
      --muted: #94a3b8;
      --blue: #60a5fa;
      --cyan: #22d3ee;
      --green: #34d399;
      --shadow: 0 24px 80px rgba(0,0,0,.45);
    }

    * { box-sizing: border-box; }

    body {
      margin: 0;
      min-height: 100vh;
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      background:
        radial-gradient(circle at 15% 10%, rgba(59,130,246,.22), transparent 28%),
        radial-gradient(circle at 85% 18%, rgba(34,211,238,.16), transparent 26%),
        radial-gradient(circle at 50% 95%, rgba(99,102,241,.18), transparent 34%),
        var(--bg);
      color: var(--text);
    }

    .shell { width: min(1180px, calc(100% - 32px)); margin: 0 auto; padding: 28px 0; }
    .nav { display: flex; justify-content: space-between; align-items: center; padding: 18px 20px; border: 1px solid var(--border); background: rgba(2, 6, 23, .72); backdrop-filter: blur(18px); border-radius: 24px; box-shadow: var(--shadow); }
    .brand { display: flex; gap: 14px; align-items: center; }
    .logo { width: 42px; height: 42px; border-radius: 14px; background: linear-gradient(135deg, var(--blue), var(--cyan)); display: grid; place-items: center; color: #03101f; font-weight: 900; box-shadow: 0 0 36px rgba(96,165,250,.35); }
    .brand h1 { font-size: 18px; margin: 0; letter-spacing: -0.03em; }
    .brand p { margin: 2px 0 0; color: var(--muted); font-size: 13px; }
    .actions { display: flex; gap: 10px; align-items: center; }
    button { border: 0; cursor: pointer; border-radius: 14px; padding: 12px 15px; font-weight: 700; color: #07111f; background: linear-gradient(135deg, #dbeafe, #67e8f9); transition: transform .12s ease, opacity .12s ease; }
    button:hover { transform: translateY(-1px); }
    button:disabled { opacity: .55; cursor: not-allowed; }
    .ghost { background: rgba(15, 23, 42, .9); color: var(--text); border: 1px solid var(--border); }
    .hero { padding: 44px 10px 26px; text-align: center; }
    .hero .kicker { display: inline-flex; gap: 8px; align-items: center; border: 1px solid rgba(96,165,250,.28); color: #bfdbfe; padding: 8px 12px; border-radius: 999px; background: rgba(37, 99, 235, .12); font-size: 13px; margin-bottom: 18px; }
    .hero h2 { font-size: clamp(34px, 6vw, 68px); line-height: .96; margin: 0 auto; max-width: 900px; letter-spacing: -0.07em; font-weight: 850; }
    .hero h2 span { color: transparent; background: linear-gradient(135deg, #93c5fd, #22d3ee); -webkit-background-clip: text; background-clip: text; }
    .hero p { max-width: 760px; margin: 22px auto 0; color: var(--muted); font-size: 17px; line-height: 1.6; }
    .grid { display: grid; grid-template-columns: 1.1fr .9fr; gap: 18px; align-items: stretch; }
    .card { border: 1px solid var(--border); background: var(--panel); backdrop-filter: blur(18px); border-radius: 28px; box-shadow: var(--shadow); overflow: hidden; }
    .chat { min-height: 610px; display: flex; flex-direction: column; }
    .chat-header { padding: 18px 20px; border-bottom: 1px solid var(--border); display: flex; justify-content: space-between; gap: 12px; align-items: center; background: rgba(15, 23, 42, .72); }
    .chat-header strong { font-size: 15px; }
    .status { display: inline-flex; gap: 8px; align-items: center; color: var(--muted); font-size: 13px; }
    .dot { width: 9px; height: 9px; border-radius: 50%; background: var(--green); box-shadow: 0 0 16px rgba(52,211,153,.8); }
    .messages { flex: 1; overflow-y: auto; padding: 22px; display: flex; flex-direction: column; gap: 14px; }
    .msg { max-width: 88%; padding: 14px 16px; border-radius: 18px; line-height: 1.55; white-space: pre-wrap; font-size: 15px; }
    .user { align-self: flex-end; background: linear-gradient(135deg, rgba(96,165,250,.95), rgba(34,211,238,.9)); color: #04111f; border-bottom-right-radius: 6px; font-weight: 650; }
    .bot { align-self: flex-start; background: rgba(15, 23, 42, .95); border: 1px solid var(--border); border-bottom-left-radius: 6px; }
    .meta { margin-top: 10px; padding-top: 10px; border-top: 1px solid rgba(148,163,184,.16); color: var(--muted); font-size: 12px; }
    .composer { padding: 18px; border-top: 1px solid var(--border); display: flex; gap: 12px; background: rgba(2, 6, 23, .38); }
    textarea { flex: 1; min-height: 54px; max-height: 140px; resize: vertical; border: 1px solid var(--border); border-radius: 18px; background: rgba(2, 6, 23, .74); color: var(--text); padding: 14px 15px; outline: none; font: inherit; }
    textarea:focus { border-color: rgba(96,165,250,.65); box-shadow: 0 0 0 4px rgba(96,165,250,.12); }
    .side { padding: 20px; display: flex; flex-direction: column; gap: 16px; }
    .side h3 { margin: 0; font-size: 18px; letter-spacing: -0.03em; }
    .side p { color: var(--muted); margin: 4px 0 0; line-height: 1.5; font-size: 14px; }
    .pill-grid, .examples { display: grid; gap: 10px; }
    .pill { padding: 14px; border: 1px solid var(--border); background: var(--panel2); border-radius: 18px; }
    .pill strong { display: block; margin-bottom: 5px; color: #dbeafe; }
    .example { text-align: left; background: rgba(2, 6, 23, .62); border: 1px solid var(--border); color: var(--text); padding: 12px 13px; border-radius: 16px; font-weight: 600; }
    .footer-note { color: var(--muted); font-size: 12px; line-height: 1.45; padding: 0 4px; }
    @media (max-width: 920px) { .grid { grid-template-columns: 1fr; } .actions { display: none; } .msg { max-width: 96%; } }
  </style>
</head>
<body>
  <main class="shell">
    <nav class="nav">
      <div class="brand">
        <div class="logo">M1</div>
        <div>
          <h1>Asistente IA · Módulo 1</h1>
          <p>RAG + búsqueda web + contexto fijo + memoria de sesión</p>
        </div>
      </div>
      <div class="actions">
        <button class="ghost" onclick="health()">Health</button>
        <button class="ghost" onclick="clearMemory()">Limpiar memoria</button>
        <button onclick="rebuild()">Reconstruir índice</button>
      </div>
    </nav>

    <section class="hero">
      <div class="kicker">⚡ Demo real · LangChain + OpenAI + Chroma</div>
      <h2>Preguntas del módulo, <span>respuestas con criterio</span></h2>
      <p>
        Este asistente usa contexto fijo del curso, memoria conversacional, RAG sobre las clases
        y búsqueda web cuando hace falta. Además recomienda cuándo usar API simple, chain, RAG,
        tool, agente o nada de IA.
      </p>
    </section>

    <section class="grid">
      <div class="card chat">
        <div class="chat-header">
          <strong>Chat académico</strong>
          <span class="status"><span class="dot"></span><span id="status">Listo</span></span>
        </div>

        <div class="messages" id="messages">
          <div class="msg bot">
            Hola. Soy el asistente del Módulo 1. Tengo contexto fijo del curso, memoria de sesión, RAG sobre las clases y búsqueda web.
            <div class="meta">Modo inicial: esperando pregunta · Memoria activa</div>
          </div>
        </div>

        <div class="composer">
          <textarea id="question" placeholder="Ej: Tengo PDFs de importación y quiero pasar datos a Excel con alertas. ¿Qué arquitectura conviene?"></textarea>
          <button id="sendBtn" onclick="ask()">Enviar</button>
        </div>
      </div>

      <aside class="card side">
        <div>
          <h3>Qué puede hacer</h3>
          <p>Diseñado para clase en vivo: útil, visual y defendible.</p>
        </div>

        <div class="pill-grid">
          <div class="pill"><strong>Contexto fijo</strong>Conoce profesor, módulo, programa, institución y reglas estables.</div>
          <div class="pill"><strong>Memoria de sesión</strong>Recuerda los últimos turnos mientras no reinicies el servidor.</div>
          <div class="pill"><strong>RAG del módulo</strong>Responde desde PPTs de clase cargados en <code>data/</code>.</div>
          <div class="pill"><strong>Web search</strong>Usa internet cuando la pregunta pide ejemplos externos o información vigente.</div>
          <div class="pill"><strong>Decisión de arquitectura</strong>Recomienda API simple, chain, RAG, tool, agente, dashboard o no IA.</div>
        </div>

        <div>
          <h3>Preguntas demo</h3>
          <p>Haz clic para probar.</p>
        </div>

        <div class="examples">
          <button class="example" onclick="fillExample('¿Quién es el profesor del módulo?')">Profesor del módulo</button>
          <button class="example" onclick="fillExample('¿Qué es Node y por qué se usa en una demo de APIs?')">¿Qué es Node?</button>
          <button class="example" onclick="fillExample('Dame ejemplos de secret managers usados en empresas.')">Ejemplos de secret managers</button>
          <button class="example" onclick="fillExample('Tengo PDFs de importación marítima, quiero extraer datos a Excel y generar alertas por cambios de fecha. ¿Uso API simple, chain, RAG, tool, agente o no IA?')">Caso importaciones PDF → Excel</button>
          <button class="example" onclick="fillExample('¿Y en ese caso por qué no usarías un agente al inicio?')">Prueba memoria</button>
          <button class="example" onclick="fillExample('¿Cuándo conviene usar RAG y cuándo basta una chain simple?')">RAG vs chain</button>
          <button class="example" onclick="fillExample('Tengo un dashboard de KPIs de proyectos de construcción. ¿Dónde sí aporta IA y dónde no?')">KPIs construcción</button>
        </div>

        <div class="footer-note">
          Consejo para clase: pregunta primero quién es el profesor para mostrar contexto fijo; luego pregunta “¿y en ese caso?” para mostrar memoria.
        </div>
      </aside>
    </section>
  </main>

  <script>
    const messages = document.getElementById("messages");
    const questionBox = document.getElementById("question");
    const statusEl = document.getElementById("status");
    const sendBtn = document.getElementById("sendBtn");

    const SESSION_ID_KEY = "modulo1_agent_session_id";
    let sessionId = localStorage.getItem(SESSION_ID_KEY);

    if (!sessionId) {
      sessionId = "session_" + crypto.randomUUID();
      localStorage.setItem(SESSION_ID_KEY, sessionId);
    }

    function addMessage(text, role, meta = "") {
      const div = document.createElement("div");
      div.className = `msg ${role}`;
      div.textContent = text;

      if (meta) {
        const m = document.createElement("div");
        m.className = "meta";
        m.textContent = meta;
        div.appendChild(m);
      }

      messages.appendChild(div);
      messages.scrollTop = messages.scrollHeight;
    }

    function fillExample(text) {
      questionBox.value = text;
      questionBox.focus();
    }

    async function ask() {
      const question = questionBox.value.trim();
      if (!question) return;

      addMessage(question, "user");
      questionBox.value = "";
      statusEl.textContent = "Pensando...";
      sendBtn.disabled = true;

      try {
        const res = await fetch("/api/ask", {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: JSON.stringify({question, session_id: sessionId})
        });

        const data = await res.json();
        const sources = data.sources && data.sources.length
          ? "Fuentes: " + data.sources.join(" · ")
          : "Sin fuentes recuperadas";

        addMessage(data.answer, "bot", `Modo: ${data.mode} · ${sources}`);
      } catch (err) {
        addMessage("Error llamando al backend: " + err.message, "bot", "Error");
      } finally {
        statusEl.textContent = "Listo";
        sendBtn.disabled = false;
      }
    }

    async function rebuild() {
      statusEl.textContent = "Reconstruyendo índice...";
      try {
        const res = await fetch("/api/rebuild", {method: "POST"});
        const data = await res.json();
        addMessage(data.message, "bot", "Índice actualizado");
      } catch (err) {
        addMessage("Error reconstruyendo índice: " + err.message, "bot", "Error");
      } finally {
        statusEl.textContent = "Listo";
      }
    }

    async function clearMemory() {
      statusEl.textContent = "Limpiando memoria...";
      try {
        const res = await fetch("/api/memory/clear", {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: JSON.stringify({session_id: sessionId})
        });
        const data = await res.json();
        addMessage(data.message, "bot", "Memoria");
      } catch (err) {
        addMessage("Error limpiando memoria: " + err.message, "bot", "Error");
      } finally {
        statusEl.textContent = "Listo";
      }
    }

    async function health() {
      try {
        const res = await fetch("/api/health");
        const data = await res.json();
        addMessage(JSON.stringify(data, null, 2), "bot", "Health check");
      } catch (err) {
        addMessage("Error en health: " + err.message, "bot", "Error");
      }
    }

    questionBox.addEventListener("keydown", function(e) {
      if (e.key === "Enter" && (e.ctrlKey || e.metaKey)) {
        ask();
      }
    });
  </script>
</body>
</html>
"""